from pptx import Presentation

# Create a new presentation
from pptx.util import Inches

prs = Presentation()

# Add a title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "ChatGPT: An Overview"
subtitle.text = "Architecture and Use Cases"

# Add a slide
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)

# Add textbox
left = top = Inches(1)
width = Inches(8)
height = Inches(5.5)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame

# Set text
tf.text = """ChatGPT is a transformer-based language model designed to facilitate conversation between humans and computer programs. It combines the power of the transformer architecture with the ability to generate natural language responses, enabling it to handle complex conversation tasks with great accuracy."""

# Add a slide
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)

# Add textbox
left = top = Inches(1)
width = Inches(8)
height = Inches(5.5)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame

# Set text
tf.text = """ChatGPT can be used in various applications such as customer service, virtual assistants, and conversational AI. By leveraging the power of natural language processing, it can accurately understand the intent of a user and generate the appropriate response. Furthermore, ChatGPT can even make suggestions to the user based on their conversation context, making the conversation more natural and efficient."""

# Save
prs.save('chatgpt.pptx')