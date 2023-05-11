import pptx
from pptx.util import Inches

# 创建一个新的PPT文件
prs = pptx.Presentation()

# 添加标题页
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "互联网金融介绍"
subtitle.text = "Python生成的PPT内容"

# 添加正文内容
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "什么是互联网金融？"
body = slide.shapes.placeholders[1]
tf = body.text_frame
tf.text = "互联网金融是指运用互联网、移动互联网、大数据、云计算、人工智能等技术手段，为金融业务提供技术支持和服务，打破传统金融机构的地域限制、时间限制、信息不对称等传统障碍，实现金融服务的普惠化、高效化、智能化。"

# 添加图片
img_path = "internet-finance.png"
slide = prs.slides.add_slide(prs.slide_layouts[1])
left = Inches(1)
top = Inches(2)
height = Inches(4)
pic = slide.shapes.add_picture(img_path, left, top, height=height)

# 保存PPT文件
prs.save("internet-finance.pptx")
